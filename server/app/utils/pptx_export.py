from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from ..models import DocumentSection, Project


def build_pptx(project: Project, sections: list[DocumentSection]) -> BytesIO:
    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = project.title
    subtitle = title_slide.placeholders[1]
    subtitle.text = project.topic

    bullet_layout = presentation.slide_layouts[1]
    for section in sorted(sections, key=lambda s: s.position):
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section.title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        body.text = ""
        for line in section.content.split("\n"):
            if line.strip():
                if not body.text:
                    body.text = line.strip()
                else:
                    body.add_paragraph().text = line.strip()

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer

